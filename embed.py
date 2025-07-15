import os
import re
import PyPDF2
from sentence_transformers import SentenceTransformer
import pinecone
from pptx import Presentation
from PIL import Image
import pytesseract
from pinecone import Pinecone, ServerlessSpec

# Initialize Pinecone
api_key = "d32b39c0-2fdb-4652-b6a3-0e1084a88325"
pc = Pinecone(api_key=api_key)
index_name = "anand-chatbot"

# Create the index if it doesn't exist
if index_name not in pc.list_indexes().names():
    pc.create_index(
        name=index_name,
        dimension=4096,
        metric='cosine',
        spec=ServerlessSpec(
            cloud='aws',
            region='us-east-1'
        )
    )
index = pc.Index(index_name)
 
# Load SentenceTransformer model
model = SentenceTransformer("nvidia/NV-Embed-v2", trust_remote_code="True")

# Directory containing files
file_directory = "C:/Users/admin/Desktop/QBOT/LLM Model Build/Database Files/"

# Define a maximum length for chunk text in metadata
MAX_CHUNK_TEXT_LENGTH = 500  # Adjust this value as needed

# Function to remove URLs
def remove_urls(text):
    url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
    return re.sub(url_pattern, '', text)

# Function to remove email addresses
def remove_emails(text):
    email_pattern = r'\S+@\S+'
    return re.sub(email_pattern, '', text)

# Function to remove page numbers (assuming "Page x of y" format)
def remove_page_numbers(text):
    page_number_pattern = r'Page \d+( of \d+)?'
    return re.sub(page_number_pattern, '', text)

# Function to remove special characters, keeping basic punctuation
def remove_special_characters(text):
    return re.sub(r'[^a-zA-Z0-9\s.,!?\'"-]', '', text)

# Function to remove extra whitespaces
def remove_extra_whitespaces(text):
    return re.sub(r'\s+', ' ', text).strip()

# Function to remove numeric only lines
def remove_numeric_only_lines(text):
    return re.sub(r'^\d+$', '', text, flags=re.MULTILINE)

# Function to clean the text
def clean_text(text):
    text = remove_urls(text)
    text = remove_emails(text)
    text = remove_page_numbers(text)
    text = remove_special_characters(text)
    text = remove_extra_whitespaces(text)
    text = remove_numeric_only_lines(text)
    return text

# Function to truncate the text to the maximum allowed length
def truncate_text(text, max_length):
    if len(text) > max_length:
        return text[:max_length] + "..."
    return text

# Function to read text from a PDF file
def read_pdf(file_path):
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in range(len(reader.pages)):
            text += reader.pages[page].extract_text()
        return text

# Updated code for extracting text from a PowerPoint file
def extract_text_from_shape(shape):
    """Extract text from a shape, including text frames and tables."""
    text = ""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text
            text += "\n"
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text + " "
            text += "\n"
    return text.strip()

def extract_text_from_ppt(ppt_file):
    """Extract text from all slides and notes in a PowerPoint file."""
    prs = Presentation(ppt_file)
    text = []

    for slide in prs.slides:
        # Extract text from slide shapes
        for shape in slide.shapes:
            extracted_text = extract_text_from_shape(shape)
            if extracted_text:
                text.append(extracted_text)

        # Extract text from notes (if any)
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                extracted_text = extract_text_from_shape(shape)
                if extracted_text:
                    text.append(extracted_text)

    return "\n".join(text)

# Function to read text from a PowerPoint file
def read_pptx(file_path):
    return extract_text_from_ppt(file_path)

# Function to read text from an image file using OCR
def read_image(file_path):
    text = pytesseract.image_to_string(Image.open(file_path))
    return text

# Function to split text into chunks
def split_text(text, chunk_size=500):
    sentences = text.split('. ')
    chunks = []
    chunk = ""
    for sentence in sentences:
        if len(chunk) + len(sentence) > chunk_size:
            chunks.append(chunk)
            chunk = sentence + ". "
        else:
            chunk += sentence + ". "
    if chunk:
        chunks.append(chunk)
    return chunks

# Iterate through the files and process them
for file_name in os.listdir(file_directory):
    file_path = os.path.join(file_directory, file_name)
    if file_name.endswith(".pdf"):
        text = read_pdf(file_path)
    elif file_name.endswith(".pptx"):
        text = read_pptx(file_path)
    elif file_name.endswith((".jpg", ".jpeg", ".png")):
        text = read_image(file_path)
    else:
        continue

    # Clean the extracted text
    cleaned_text = clean_text(text)
    chunks = split_text(cleaned_text)
    for i, chunk in enumerate(chunks):
        embedding = model.encode(chunk)
        truncated_chunk = truncate_text(chunk, MAX_CHUNK_TEXT_LENGTH)
        metadata = {
            "file_name": file_name,
            "chunk_index": i,
            "chunk_text": truncated_chunk
        }
        index.upsert([("file_chunk_" + file_name + "_" + str(i), embedding.tolist(), metadata)])